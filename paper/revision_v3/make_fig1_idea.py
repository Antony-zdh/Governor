#!/usr/bin/env python3
"""Figure 1 (idea figure) for the False Consensus paper -> editable .pptx + .pdf.

Spec: paper/revision_v3/FIG1_IDEA_SPEC.md
Panel A draws the REAL probe stream of MATH500 problem 68
(DeepSeek-R1-Distill-Qwen-7B, seed 42, dense 64-token probes).

Usage:
    python3 make_fig1_idea.py [--out DIR]
    # then: soffice --headless --convert-to pdf fig1_idea.pptx
"""
from __future__ import annotations

import argparse
import csv
import collections
import json
import os
import subprocess
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# palette (matches report/make_generalization_figs.py)
# --------------------------------------------------------------------------- #
INK        = RGBColor(0x1F, 0x23, 0x28)
MUTED      = RGBColor(0x6B, 0x72, 0x80)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

CONS_FILL  = RGBColor(0xFD, 0xE6, 0x8A)   # consensus / amber
CONS_LINE  = RGBColor(0xF5, 0x9E, 0x0B)
CONS_TEXT  = RGBColor(0xB4, 0x53, 0x09)
CONS_SOLID = RGBColor(0xB4, 0x53, 0x09)

DEER_FILL  = RGBColor(0xD1, 0xFA, 0xE5)   # DEER / emerald
DEER_LINE  = RGBColor(0x05, 0x96, 0x69)
DEER_TEXT  = RGBColor(0x06, 0x5F, 0x46)

WRONG      = RGBColor(0xDC, 0x26, 0x26)   # dominant wrong answer
WRONG_ALT  = RGBColor(0xFC, 0xA5, 0xA5)   # other wrong answers
RIGHT      = RGBColor(0x05, 0x96, 0x69)
INVALID    = RGBColor(0xE5, 0xE7, 0xEB)

PANEL_FILL = RGBColor(0xF8, 0xF9, 0xFA)
PANEL_LINE = RGBColor(0xD9, 0xDD, 0xE1)
GATE_FILL  = RGBColor(0xE5, 0xE7, 0xEB)
GATE_LINE  = RGBColor(0x9C, 0xA3, 0xAF)

FONT = "Helvetica"

# --------------------------------------------------------------------------- #
# geometry (inches)
# --------------------------------------------------------------------------- #
W, H = 7.0, 2.75
M = 0.06
BAND_H = 0.28
PANEL_Y = M
PANEL_H = H - 2 * M - BAND_H - 0.05
GUT = 0.36
AW = 1.80
BW = 2.36
AX = M
BX = AX + AW + GUT
CX = BX + BW + GUT
CW = W - M - CX

# --------------------------------------------------------------------------- #
# helpers
# --------------------------------------------------------------------------- #


def _shape(slide, kind, x, y, w, h, fill=None, line=None, lw=0.75):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(s.text_frame, m, 0)
    return s


def rect(slide, x, y, w, h, **kw):
    return _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)


def round_rect(slide, x, y, w, h, radius=0.12, **kw):
    s = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kw)
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s


def text(slide, x, y, w, h, runs, size=7.5, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
         line_spacing=0.95):
    """runs: str, or list of (str, dict-of-overrides) for mixed formatting,
    or list of lists (one per paragraph)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = anchor

    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
        items = para if isinstance(para, list) else [(para, {})]
        if isinstance(items, str):
            items = [(items, {})]
        for chunk in items:
            txt, ov = chunk if isinstance(chunk, tuple) else (chunk, {})
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.color.rgb = ov.get("color", color)
    return tb


def label_shape(shape, runs, size=7.5, color=INK, bold=False,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=False):
    tf = shape.text_frame
    tf.vertical_anchor = anchor
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try:
            p.line_spacing = 0.92
        except Exception:
            pass
        items = para if isinstance(para, list) else [(para, {})]
        if isinstance(items, str):
            items = [(items, {})]
        for chunk in items:
            txt, ov = chunk if isinstance(chunk, tuple) else (chunk, {})
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = FONT
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", italic)
            f.color.rgb = ov.get("color", color)
    return shape


def arrow(slide, x, y, w, h, color=MUTED):
    s = _shape(slide, MSO_SHAPE.RIGHT_ARROW, x, y, w, h, fill=color, line=None)
    return s


def funnel_icon(slide, x, y, w, h, fill=GATE_FILL, line=GATE_LINE):
    """A funnel: wide bowl tapering to a neck, then a short stem. Drawn as one
    freeform so nothing has to be rotated and no text sits on top of it."""
    nw = w * 0.16               # neck width
    bh = h * 0.62               # bowl height
    b = slide.shapes.build_freeform(Inches(x), Inches(y))
    b.add_line_segments([
        (Inches(x + w), Inches(y)),
        (Inches(x + w / 2 + nw / 2), Inches(y + bh)),
        (Inches(x + w / 2 + nw / 2), Inches(y + h)),
        (Inches(x + w / 2 - nw / 2), Inches(y + h)),
        (Inches(x + w / 2 - nw / 2), Inches(y + bh)),
    ], close=True)
    s = b.convert_to_shape()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def funnel(slide, x, y, w, h, narrow=0.42, fill=GATE_FILL, line=GATE_LINE):
    """Downward-narrowing trapezoid as a freeform (no rotation, so any text
    placed over it stays upright in both PowerPoint and LibreOffice)."""
    inset = w * (1 - narrow) / 2
    b = slide.shapes.build_freeform(Inches(x), Inches(y))
    b.add_line_segments([(Inches(x + w), Inches(y)),
                         (Inches(x + w - inset), Inches(y + h)),
                         (Inches(x + inset), Inches(y + h))], close=True)
    s = b.convert_to_shape()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


# --------------------------------------------------------------------------- #
# panel A data: the real probe stream
# --------------------------------------------------------------------------- #
REPO = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
ENV = os.path.join(
    REPO, "benchmark/FalseConsensus/results/governor_v2",
    "development__deepseek-ai-deepseek-r1-distill-qwen-7b__math500__seed_42")
PID = 68

# Fallback: the verified stream for MATH500 #68 (gold 46, final 46, correct).
FALLBACK = (["45"] + ["37"] * 10 + ["134", "134", "37", "32", "45", "37", "32"]
            + ["46"] * 47)


def load_stream():
    csv_path = os.path.join(ENV, "dense_simple32", "probes.csv")
    traj_path = os.path.join(ENV, "main", "traj", f"problem_{PID}.json")
    if not (os.path.exists(csv_path) and os.path.exists(traj_path)):
        sys.stderr.write("! probe bank not found, using verified fallback stream\n")
        return FALLBACK, "46", True
    rows = []
    with open(csv_path) as f:
        for r in csv.DictReader(f):
            if int(r["problem_id"]) == PID:
                rows.append((int(r["token_position"]), r["probe_answer"]))
    rows.sort()
    t = json.load(open(traj_path))
    return [a for _, a in rows], t["target"], bool(t["final_correct"])


def panel_head(slide, x, y, w, letter, title, subtitle, accent):
    text(slide, x + 0.09, y + 0.075, w - 0.18, 0.15,
         [[(letter, {"bold": True, "color": accent, "size": 8.5}),
           ("   " + title, {"bold": True, "size": 8.5})]])
    text(slide, x + 0.09, y + 0.225, w - 0.18, 0.12, subtitle, size=6.0,
         color=MUTED, italic=True)


def draw_panel_a(slide):
    round_rect(slide, AX, PANEL_Y, AW, PANEL_H, fill=PANEL_FILL, line=PANEL_LINE)
    panel_head(slide, AX, PANEL_Y, AW, "a", "Agreement ≠ terminality",
               "the real probe stream of one trajectory", CONS_TEXT)

    ix, iw = AX + 0.11, AW - 0.22

    # headline claim
    text(slide, ix, PANEL_Y + 0.40, iw, 0.42,
         [[("The answer is still moving.", {"bold": True, "size": 7.2})],
          [("Stopping on agreement destroys a recovery the",
            {"color": MUTED, "size": 6.0})],
          [("trajectory would have made on its own.",
            {"color": MUTED, "size": 6.0})]], line_spacing=1.05)

    stream, gold, ok = load_stream()
    n = len(stream)
    rx, rw = ix, iw
    ry, rh = PANEL_Y + 1.02, 0.21
    cw = rw / n

    for i, a in enumerate(stream):
        if a == gold:
            c = RIGHT
        elif not a or a.strip() in {"", "-"}:
            c = INVALID
        elif a == "37":
            c = WRONG
        else:
            c = WRONG_ALT
        rect(slide, rx + i * cw, ry, cw * 1.03, rh, fill=c, line=None)
    rect(slide, rx - 0.004, ry - 0.004, rw + 0.008, rh + 0.008, fill=None,
         line=RGBColor(0xC7, 0xCC, 0xD1), lw=0.5)

    # labels above the ribbon (two non-overlapping halves)
    text(slide, rx, ry - 0.145, rw * 0.46, 0.12,
         [[("10 ×  “37”  ✗", {"bold": True, "color": WRONG, "size": 6.2})]])
    text(slide, rx + rw * 0.48, ry - 0.145, rw * 0.52, 0.12,
         [[("“46”  ✓  to the end",
            {"bold": True, "color": DEER_TEXT, "size": 6.2})]],
         align=PP_ALIGN.RIGHT)

    # token axis
    rect(slide, rx, ry + rh + 0.03, rw, 0.007, fill=RGBColor(0xC7, 0xCC, 0xD1),
         line=None)
    text(slide, rx + rw - 0.90, ry + rh + 0.05, 0.90, 0.10,
         "4,197 generated tokens", size=5.2, color=MUTED, align=PP_ALIGN.RIGHT)

    # stop marker: a W=3, s=1.0 rule fires at the third agreeing probe
    sx = rx + 3.5 * cw
    tri = _shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, sx - 0.042, ry + rh + 0.005,
                 0.084, 0.062, fill=WRONG, line=None)
    tri.rotation = 180
    text(slide, ix, ry + rh + 0.175, iw, 0.16,
         [[("consensus fires at 256, on “37”",
            {"bold": True, "color": WRONG, "size": 6.2})]])

    # badge
    badge = round_rect(slide, ix, PANEL_Y + PANEL_H - 0.36, iw, 0.27, radius=0.30,
                       fill=RGBColor(0xFE, 0xF3, 0xC7), line=CONS_LINE)
    label_shape(badge,
                [[("harm : rescue ≈ 45 : 1", {"bold": True, "color": CONS_TEXT,
                                              "size": 7.4})]])


# --------------------------------------------------------------------------- #
def draw_panel_b(slide):
    round_rect(slide, BX, PANEL_Y, BW, PANEL_H, fill=PANEL_FILL, line=PANEL_LINE)
    panel_head(slide, BX, PANEL_Y, BW, "b", "3,520 rules, gates fixed in advance",
               "preregistered sweep on frozen trajectories", CONS_TEXT)

    ix, iw = BX + 0.11, BW - 0.22

    # --- row 1: the evidence base ----------------------------------------- #
    y1 = PANEL_Y + 0.44
    r1 = round_rect(slide, ix, y1, iw, 0.26, radius=0.14, fill=WHITE, line=PANEL_LINE)
    label_shape(r1, [[("18 frozen environments", {"bold": True, "size": 7.0})],
                     [("2 models × 3 benchmarks × 3 seeds",
                       {"color": MUTED, "size": 5.6})]])

    # --- row 2: the rule space, with axis labels aligned to the cells ------ #
    # the grid + its axis labels + the "3,520" caption are treated as one block
    # and centred, so the panel does not read as left-weighted.
    y2 = y1 + 0.42
    cell, gap = 0.052, 0.013
    step = cell + gap
    gw = 8 * step - gap
    lab_w, cap_w = 0.24, 0.86          # s-labels on the left, caption on the right
    block_w = lab_w + gw + 0.12 + cap_w
    bx0 = ix + (iw - block_w) / 2
    gx0 = bx0 + lab_w
    for r in range(3):
        for c in range(8):
            skip = (r > 0 and c == 0)      # W=1, s != 1.0 is behaviourally redundant
            rect(slide, gx0 + c * step, y2 + r * step, cell, cell,
                 fill=INVALID if skip else CONS_FILL,
                 line=PANEL_LINE if skip else CONS_LINE, lw=0.5)
    for r, lab in enumerate(("1.0", "0.8", "0.6")):
        text(slide, bx0, y2 + r * step + cell / 2 - 0.035, lab_w - 0.035, 0.07,
             lab, size=5.0, color=MUTED, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
    text(slide, bx0, y2 - 0.10, lab_w - 0.035, 0.08, "s", size=5.4, color=MUTED,
         align=PP_ALIGN.RIGHT, italic=True)
    for c, lab in ((0, "1"), (7, "30")):
        text(slide, gx0 + c * step - 0.06, y2 + 3 * step + 0.005, cell + 0.12, 0.08,
             lab, size=5.0, color=MUTED, align=PP_ALIGN.CENTER)
    text(slide, gx0 + gw / 2 - 0.15, y2 + 3 * step + 0.005, 0.30, 0.08, "W",
         size=5.4, color=MUTED, align=PP_ALIGN.CENTER, italic=True)
    text(slide, gx0 + gw + 0.12, y2 + 0.005, cap_w, 0.24,
         [[("3,520", {"bold": True, "color": CONS_TEXT, "size": 9.8})],
          [("consensus rules", {"color": CONS_TEXT, "size": 6.2})]],
         align=PP_ALIGN.CENTER, line_spacing=1.02)

    # --- row 3: the three preregistered gates ------------------------------ #
    y3 = y2 + 3 * step + 0.30
    grow_w = 1.46
    gx = ix + (iw - grow_w) / 2
    text(slide, gx, y3 - 0.115, grow_w, 0.09,
         "three preregistered gates", size=5.6, color=MUTED, italic=True)
    for k, (lab, cap, floor, col) in enumerate(
            (("C", "1.0", "10", RGBColor(0x16, 0xA3, 0x4A)),
             ("B", "2.0", "20", RGBColor(0x25, 0x63, 0xEB)),
             ("T", "3.5", "30", RGBColor(0x93, 0x33, 0xEA)))):
        yy = y3 + k * 0.145
        d = _shape(slide, MSO_SHAPE.OVAL, gx, yy + 0.019, 0.090, 0.090,
                   fill=col, line=None)
        label_shape(d, [[(lab, {"bold": True, "size": 5.4, "color": WHITE})]])
        text(slide, gx + 0.125, yy, grow_w - 0.125, 0.12,
             [[("drop ≤ ", {"size": 6.2, "color": MUTED}),
               (f"{cap} pp", {"size": 6.2}),
               ("    saving ≥ ", {"size": 6.2, "color": MUTED}),
               (f"{floor}%", {"size": 6.2})]], anchor=MSO_ANCHOR.MIDDLE)

    # --- row 4: the verdict ------------------------------------------------ #
    y4 = y3 + 0.56
    v = round_rect(slide, ix + 0.18, y4, iw - 0.36, 0.34, radius=0.16,
                   fill=CONS_SOLID, line=None)
    label_shape(v, [[("0 / 3,520", {"bold": True, "size": 13.0, "color": WHITE}),
                     ("   clear any gate", {"size": 7.4, "color": WHITE})]])


def draw_panel_c(slide):
    round_rect(slide, CX, PANEL_Y, CW, PANEL_H, fill=PANEL_FILL, line=PANEL_LINE)
    panel_head(slide, CX, PANEL_Y, CW, "c", "The failure is the signal",
               "identical pipeline, gates and accounting", DEER_TEXT)

    ix, iw = CX + 0.11, CW - 0.22
    hw = (iw - 0.08) / 2

    # two candidate signals
    y1 = PANEL_Y + 0.50
    s1 = round_rect(slide, ix, y1, hw, 0.36, radius=0.12,
                    fill=RGBColor(0xFE, 0xF3, 0xC7), line=CONS_LINE)
    label_shape(s1, [[("consensus", {"bold": True, "size": 7.0, "color": CONS_TEXT})],
                     [("do recent probes", {"size": 5.4, "color": MUTED})],
                     [("agree?", {"size": 5.4, "color": MUTED})]])
    s2 = round_rect(slide, ix + hw + 0.08, y1, hw, 0.36, radius=0.12,
                    fill=DEER_FILL, line=DEER_LINE)
    label_shape(s2, [[("DEER", {"bold": True, "size": 7.0, "color": DEER_TEXT})],
                     [("confident at a", {"size": 5.4, "color": MUTED})],
                     [("reasoning boundary?", {"size": 5.4, "color": MUTED})]])

    # the shared pipeline
    y2 = y1 + 0.46
    for cx0 in (ix + hw / 2, ix + hw + 0.08 + hw / 2):
        _shape(slide, MSO_SHAPE.DOWN_ARROW, cx0 - 0.028, y2 - 0.062, 0.056, 0.056,
               fill=MUTED, line=None)
    pipe = round_rect(slide, ix, y2, iw, 0.23, radius=0.30,
                      fill=RGBColor(0x6B, 0x72, 0x80), line=None)
    label_shape(pipe, [[("same sweep · same gates · same accounting",
                         {"size": 6.0, "color": WHITE, "bold": True})]])

    # verdicts
    y3 = y2 + 0.42
    v1 = round_rect(slide, ix, y3, iw, 0.30, radius=0.14, fill=WHITE, line=CONS_LINE)
    label_shape(v1, [[("✗  consensus", {"bold": True, "size": 7.6, "color": CONS_TEXT}),
                      ("    0 of 3 gates", {"size": 7.2, "color": WRONG})]])
    v2 = round_rect(slide, ix, y3 + 0.38, iw, 0.38, radius=0.14, fill=DEER_FILL,
                    line=DEER_LINE)
    label_shape(v2, [[("✓  DEER", {"bold": True, "size": 7.6, "color": DEER_TEXT}),
                      ("    3 of 3 gates", {"bold": True, "size": 7.2,
                                            "color": DEER_TEXT})],
                     [("−0.3 pp accuracy @ 28% tokens saved",
                       {"size": 6.2, "color": DEER_TEXT})]])


def draw_connectors(slide):
    """Each gutter carries one question: a soft pill with the text, and a
    rightward arrow beneath it, both centred in the gutter, so the eye reads
    panel -> question -> next panel in the direction of the argument."""
    for gx0, lines in ((AX + AW, ("can any", "rule be", "safe and", "saving?")),
                       (BX + BW, ("is early", "exit itself", "impossible?"))):
        cx = gx0 + GUT / 2
        ph = 0.098 * len(lines) + 0.09
        top = PANEL_Y + PANEL_H * 0.42 - ph / 2
        pw = GUT + 0.08          # allowed to reach a little into both panels
        pill = round_rect(slide, cx - pw / 2, top, pw, ph, radius=0.22,
                          fill=WHITE, line=PANEL_LINE)
        label_shape(pill, [[(l, {"size": 5.0, "color": MUTED, "italic": True})]
                           for l in lines])
        _shape(slide, MSO_SHAPE.RIGHT_ARROW, cx - 0.075, top + ph + 0.06,
               0.15, 0.075, fill=RGBColor(0xC7, 0xCC, 0xD1), line=None)


def draw_band(slide):
    y = H - M - BAND_H
    b = round_rect(slide, M, y, W - 2 * M, BAND_H, radius=0.10, fill=INK, line=None)
    label_shape(b, [[("Early exit is possible — the consensus ",
                      {"size": 8.5, "color": WHITE}),
                     ("signal", {"size": 8.5, "color": WHITE, "bold": True,
                                 "italic": True}),
                     (" is what cannot make it safe.",
                      {"size": 8.5, "color": WHITE})]])


# --------------------------------------------------------------------------- #
def build(out_dir):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    draw_panel_a(slide)
    draw_panel_b(slide)
    draw_panel_c(slide)
    draw_connectors(slide)
    draw_band(slide)

    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, "fig1_idea.pptx")
    prs.save(pptx_path)
    print("wrote", pptx_path)
    return pptx_path


def to_pdf(pptx_path):
    out_dir = os.path.dirname(pptx_path)
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                    "--outdir", out_dir, pptx_path], check=True,
                   stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    pdf = pptx_path.replace(".pptx", ".pdf")
    print("wrote", pdf)
    return pdf


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(REPO, "paper", "figures", "gen"))
    ap.add_argument("--no-pdf", action="store_true")
    a = ap.parse_args()
    p = build(a.out)
    if not a.no_pdf:
        to_pdf(p)
