#!/usr/bin/env python3
"""Figure 1, variant B: the five-stage pipeline from paper/revision_v3/
figure_prompts.md (issue 8), matching the five beats of CORE_PAPER_FLOW.md.

Variant A (make_fig1_idea.py) is the tighter three-panel argument; this one
follows the handoff spec literally: phenomenon -> the question -> sweep+gates
-> same gates with a different signal -> mechanism.

The mini Pareto insets in stages 3 and 4 plot REAL dev-split points (subsampled
from the committed sweep bank), drawn as native PowerPoint shapes so the whole
figure stays editable.

Usage: python3 make_fig1_idea_b.py [--out DIR]
"""
from __future__ import annotations

import argparse
import gzip
import json
import os
import random
import statistics as st
import subprocess
import sys
from collections import defaultdict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from make_fig1_idea import (  # noqa: E402
    CONS_FILL, CONS_LINE, CONS_SOLID, CONS_TEXT, DEER_FILL, DEER_LINE, DEER_TEXT,
    GATE_FILL, GATE_LINE, INK, INVALID, MUTED, PANEL_FILL, PANEL_LINE, REPO,
    RIGHT, WHITE, WRONG, WRONG_ALT, _shape, label_shape, rect, round_rect, text,
)

W, H = 7.0, 2.30
M = 0.05
BAND_H = 0.26
PY = M
PH = H - 2 * M - BAND_H - 0.04
GUT = 0.085
SW = (W - 2 * M - 4 * GUT) / 5          # stage width

BANK = os.path.join(REPO, "benchmark/FalseConsensus/results/governor_v2_ws_sweep")
SELB = {"math500": 16384, "amc23": 16384, "aime24": 32768}
DEVID = {"deepseek-ai/DeepSeek-R1-Distill-Qwen-7B", "Qwen/Qwen3-8B"}


# --------------------------------------------------------------------------- #
def load_points(n=150):
    """(drop, saving) for the consensus cloud and the DEER frontier, dev split."""
    try:
        by = defaultdict(list)
        with gzip.open(os.path.join(BANK, "dev/consensus_dev_train.jsonl.gz"), "rt") as f:
            for line in f:
                r = json.loads(line)
                if r["split"] != "dev" or r["model"] not in DEVID:
                    continue
                by[r["rule_id"]].append(r)
        cons = [(st.fmean(x["accuracy_drop_pp"] for x in v),
                 st.fmean(x["saving_fraction"] for x in v) * 100.0)
                for v in by.values()]
        dby = defaultdict(list)
        with gzip.open(os.path.join(BANK, "deer/deer_threshold_sweep.jsonl.gz"), "rt") as f:
            for line in f:
                r = json.loads(line)
                if r["split"] == "dev":
                    dby[float(r["threshold"])].append(r)
        deer = sorted((st.fmean(x["accuracy_drop_pp"] for x in v),
                       st.fmean(x["saving_fraction"] for x in v) * 100.0)
                      for v in dby.values())
    except FileNotFoundError:
        sys.stderr.write("! sweep bank not found; using a schematic cloud\n")
        rnd = random.Random(0)
        cons = [(rnd.uniform(1, 28), rnd.uniform(5, 90)) for _ in range(n)]
        deer = [(0.3, 28.2), (1.0, 29.6), (2.75, 31.9), (6.0, 40.0)]
    rnd = random.Random(7)
    if len(cons) > n:
        cons = rnd.sample(cons, n)
    return cons, deer


def mini_pareto(slide, x, y, w, h, cons, deer, *, show_deer, gate=(1.0, 10.0)):
    """A small drop-vs-saving panel drawn from native shapes."""
    xmax, ymax = 30.0, 100.0
    rect(slide, x, y, w, h, fill=WHITE, line=PANEL_LINE, lw=0.5)

    def px(d):
        return x + max(0.0, min(d, xmax)) / xmax * w

    def py(s):
        return y + h - max(0.0, min(s, ymax)) / ymax * h

    # the safe-and-saving gate box, in the top-left corner
    gx, gy = px(gate[0]), py(ymax)
    rect(slide, x, gy, gx - x, py(gate[1]) - gy,
         fill=RGBColor(0xDC, 0xFC, 0xE7), line=RGBColor(0x16, 0xA3, 0x4A), lw=0.6)

    for d, s in cons:
        r = 0.014
        _shape(slide, MSO_SHAPE.OVAL, px(d) - r, py(s) - r, 2 * r, 2 * r,
               fill=RGBColor(0xB9, 0xBE, 0xC7), line=None)
    if show_deer:
        for d, s in deer:
            r = 0.028
            _shape(slide, MSO_SHAPE.OVAL, px(d) - r, py(s) - r, 2 * r, 2 * r,
                   fill=DEER_LINE, line=None)
    text(slide, x + 0.02, gy + 0.01, w * 0.7, 0.10,
         [[("safe &", {"size": 4.0, "color": RGBColor(0x15, 0x80, 0x3D)})],
          [("saving", {"size": 4.0, "color": RGBColor(0x15, 0x80, 0x3D)})]],
         line_spacing=1.05)
    return px, py


# --------------------------------------------------------------------------- #
def stage(slide, i, letter, title, caption, accent):
    x = M + i * (SW + GUT)
    round_rect(slide, x, PY, SW, PH, fill=PANEL_FILL, line=PANEL_LINE)
    text(slide, x + 0.07, PY + 0.06, SW - 0.14, 0.30,
         [[(letter, {"bold": True, "color": accent, "size": 7.2}),
           ("  " + title, {"bold": True, "size": 7.2})]], line_spacing=1.05)
    text(slide, x + 0.07, PY + PH - 0.30, SW - 0.14, 0.30,
         [[(l, {"size": 5.6, "color": MUTED, "italic": True})]
          for l in caption.split("\n")], line_spacing=1.08)
    return x


def build(out_dir):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cons, deer = load_points()

    # ---- 1. phenomenon ---------------------------------------------------- #
    x = stage(slide, 0, "1", "Agreement is not\nterminality",
              "Probes agree while the\nanswer is still moving.", CONS_TEXT)
    ix, iw = x + 0.08, SW - 0.16
    ry = PY + 0.60
    stream = ["B", "7", "12", "12", "12", "9", "7"]
    cw = iw / len(stream)
    for k, a in enumerate(stream):
        wrong = a != "7"
        c = RGBColor(0xFE, 0xE2, 0xE2) if wrong else RGBColor(0xD1, 0xFA, 0xE5)
        e = WRONG if wrong else DEER_LINE
        b = round_rect(slide, ix + k * cw + 0.004, ry, cw - 0.008, 0.15,
                       radius=0.25, fill=c, line=e, lw=0.5)
        label_shape(b, [[(a, {"size": 5.0, "color": WRONG if wrong else DEER_TEXT})]])
    rect(slide, ix, ry + 0.185, iw, 0.02, fill=RGBColor(0xD1, 0xD5, 0xDB), line=None)
    _shape(slide, MSO_SHAPE.RIGHT_ARROW, ix + iw - 0.06, ry + 0.175, 0.06, 0.04,
           fill=RGBColor(0xD1, 0xD5, 0xDB), line=None)
    tri = _shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, ix + 2.5 * cw - 0.035,
                 ry + 0.215, 0.07, 0.055, fill=WRONG, line=None)
    tri.rotation = 180
    text(slide, ix, ry + 0.275, iw, 0.24,
         [[("stop here?", {"bold": True, "size": 5.6, "color": WRONG})],
          [("three probes agree on 12", {"size": 5.0, "color": MUTED})]],
         align=PP_ALIGN.CENTER, line_spacing=1.1)
    text(slide, ix, ry - 0.13, iw, 0.11,
         [[("true answer: ", {"size": 5.0, "color": MUTED}),
           ("7", {"size": 5.8, "color": DEER_TEXT, "bold": True})]],
         align=PP_ALIGN.CENTER)

    # ---- 2. the question -------------------------------------------------- #
    x = stage(slide, 1, "2", "So we searched the\nwhole rule space",
              "Can any consensus rule\nbe safe and saving?", CONS_TEXT)
    ix, iw = x + 0.08, SW - 0.16
    cell, gap = 0.052, 0.013
    gw = 8 * (cell + gap) - gap
    gx0 = ix + (iw - gw) / 2 + 0.05
    gy0 = PY + 0.62
    for r in range(3):
        for c in range(8):
            skip = (r > 0 and c == 0)
            rect(slide, gx0 + c * (cell + gap), gy0 + r * (cell + gap), cell, cell,
                 fill=INVALID if skip else CONS_FILL,
                 line=PANEL_LINE if skip else CONS_LINE, lw=0.5)
    text(slide, ix - 0.04, gy0 - 0.005, gx0 - ix, 0.22,
         [[("s 1.0", {"size": 4.6, "color": MUTED})],
          [("0.8", {"size": 4.6, "color": MUTED})],
          [("0.6", {"size": 4.6, "color": MUTED})]], align=PP_ALIGN.RIGHT,
         line_spacing=1.22)
    text(slide, gx0, gy0 + 3 * (cell + gap) + 0.01, gw + 0.3, 0.10,
         [[("W = 1 → 30", {"size": 4.6, "color": MUTED})]])
    text(slide, ix, gy0 + 0.29, iw, 0.30,
         [[("3,520 rules", {"bold": True, "size": 8.0, "color": CONS_TEXT})],
          [("window size × share threshold,", {"size": 5.0, "color": MUTED})],
          [("× schedule / maturity / validity", {"size": 5.0, "color": MUTED})]],
         align=PP_ALIGN.CENTER, line_spacing=1.12)

    # ---- 3. the gates say no ---------------------------------------------- #
    x = stage(slide, 2, "3", "Not one clears a\npreregistered gate",
              "Safe forces zero saving;\nsaving forces a real loss.", CONS_TEXT)
    ix, iw = x + 0.08, SW - 0.16
    mini_pareto(slide, ix + 0.10, PY + 0.60, iw - 0.20, 0.44, cons, deer,
                show_deer=False)
    text(slide, ix, PY + 1.055, iw, 0.11,
         [[("gates fixed before the sweep ran", {"size": 4.7, "color": MUTED,
                                                 "italic": True})]],
         align=PP_ALIGN.CENTER)
    v = round_rect(slide, ix + 0.06, PY + 1.17, iw - 0.12, 0.25, radius=0.16,
                   fill=CONS_SOLID, line=None)
    label_shape(v, [[("0 / 3,520", {"bold": True, "size": 10.0, "color": WHITE})]])

    # ---- 4. same gates, different signal ---------------------------------- #
    x = stage(slide, 3, "4", "A non-consensus\nsignal clears them",
              "The failure is the signal,\nnot early exit.", DEER_TEXT)
    ix, iw = x + 0.08, SW - 0.16
    mini_pareto(slide, ix + 0.10, PY + 0.60, iw - 0.20, 0.44, cons, deer,
                show_deer=True)
    text(slide, ix, PY + 1.055, iw, 0.11,
         [[("identical pipeline and gates", {"size": 4.7, "color": MUTED,
                                             "italic": True})]],
         align=PP_ALIGN.CENTER)
    v = round_rect(slide, ix + 0.06, PY + 1.17, iw - 0.12, 0.25, radius=0.16,
                   fill=DEER_LINE, line=None)
    label_shape(v, [[("DEER  3 / 3", {"bold": True, "size": 8.4, "color": WHITE})],
                    [("−0.3 pp @ 28% saved", {"size": 5.2, "color": WHITE})]])

    # ---- 5. mechanism ----------------------------------------------------- #
    x = stage(slide, 4, "5", "Why the signal\nfails", "Forced probing masks the\n"
              "model's true estimate.", CONS_TEXT)
    ix, iw = x + 0.08, SW - 0.16
    bub = _shape(slide, MSO_SHAPE.CLOUD_CALLOUT, ix, PY + 0.54, iw, 0.38,
                 fill=WHITE, line=PANEL_LINE)
    label_shape(bub, [[("“…not sure yet”", {"size": 5.6, "color": MUTED,
                                            "italic": True})]])
    text(slide, ix, PY + 0.95, iw, 0.12,
         [[("but the probe forces an answer", {"size": 5.0, "color": MUTED})]],
         align=PP_ALIGN.CENTER)
    b = round_rect(slide, ix + iw / 2 - 0.20, PY + 1.09, 0.40, 0.19, radius=0.14,
                   fill=RGBColor(0xFE, 0xE2, 0xE2), line=WRONG)
    label_shape(b, [[("12", {"bold": True, "size": 7.0, "color": WRONG})]])
    text(slide, ix, PY + 1.31, iw, 0.24,
         [[("a placeholder, repeated", {"size": 5.2, "color": WRONG})],
          [("→ read as consensus", {"size": 5.2, "color": WRONG, "bold": True})]],
         align=PP_ALIGN.CENTER, line_spacing=1.1)

    # ---- connectors + band ------------------------------------------------ #
    for i in range(4):
        gx0 = M + i * (SW + GUT) + SW
        _shape(slide, MSO_SHAPE.RIGHT_ARROW, gx0 + 0.012, PY + PH * 0.47,
               GUT - 0.024, 0.07, fill=MUTED, line=None)

    y = H - M - BAND_H
    band = round_rect(slide, M, y, W - 2 * M, BAND_H, radius=0.10, fill=INK,
                      line=None)
    label_shape(band, [[("Early exit is possible — the consensus ",
                         {"size": 8.2, "color": WHITE}),
                        ("signal", {"size": 8.2, "color": WHITE, "bold": True,
                                    "italic": True}),
                        (" is what cannot make it safe.",
                         {"size": 8.2, "color": WHITE})]])

    os.makedirs(out_dir, exist_ok=True)
    p = os.path.join(out_dir, "fig1_idea_b.pptx")
    prs.save(p)
    print("wrote", p)
    return p


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=os.path.join(REPO, "paper", "figures", "gen"))
    a = ap.parse_args()
    p = build(a.out)
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir",
                    a.out, p], check=True, stdout=subprocess.DEVNULL,
                   stderr=subprocess.DEVNULL)
    print("wrote", p.replace(".pptx", ".pdf"))
